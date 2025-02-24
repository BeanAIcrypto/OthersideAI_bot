import logging
import os
import zipfile
import py7zr
import mimetypes
import tempfile
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from typing import List, Optional

logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_path: str) -> str:
    """Извлекает текст из PDF файла."""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            text = "".join(page.extract_text() or "" for page in reader.pages)
        return text
    except Exception as e:
        logger.error(f"Ошибка при обработке PDF {file_path}: {e}")
        raise ValueError("Ошибка при извлечении текста из PDF.")

def extract_text_from_docx(file_path: str) -> str:
    """Извлекает текст из DOCX файла."""
    try:
        doc: Document = Document(file_path)
        paragraphs: List[str] = [para.text for para in doc.paragraphs]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.error(f"Ошибка при обработке DOCX файла {file_path}: {e}")
        raise Exception("Ошибка при извлечении текста из DOCX.")

def extract_text_from_excel(file_path: str) -> Optional[str]:
    """Извлекает текст из файла Excel (XLSX)."""
    try:
        df = pd.read_excel(file_path)
        df = df.loc[:, ~df.columns.str.contains("^Unnamed")]
        df = df.dropna(how="all").dropna(axis=1, how="all")
        if df.empty:
            raise ValueError(f"Файл {file_path} пустой.")
        return df.to_markdown(index=False)
    except Exception as e:
        logger.error(f"Ошибка обработки файла Excel {file_path}: {e}")
        raise Exception("Ошибка при обработке Excel файла.")

def extract_text_from_presentation(file_path: str) -> Optional[str]:
    """Извлекает текст из файла PowerPoint (PPTX)."""
    try:
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text.strip())
        if not text:
            raise ValueError(f"Файл {file_path} не содержит текста.")
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Ошибка при обработке PPTX файла {file_path}: {e}")
        raise Exception("Ошибка при обработке презентации.")

def extract_text_from_zip(file_path: str) -> Optional[str]:
    """Извлекает текст из ZIP-архива с учетом MIME-типа."""
    try:
        extracted_text: List[str] = []
        with tempfile.TemporaryDirectory() as tmpdir:
            with zipfile.ZipFile(file_path, "r") as archive:
                if not archive.namelist():
                    raise ValueError(f"Архив {file_path} пустой.")
                for file in archive.namelist():
                    mime_type = mimetypes.guess_type(file)[0]
                    if mime_type in text_extraction_from_a_document:
                        extracted_file_path = archive.extract(file, path=tmpdir)
                        text_extraction_function = text_extraction_from_a_document[mime_type]
                        text_document = text_extraction_function(extracted_file_path)
                        if text_document:
                            extracted_text.append(text_document)
                        else:
                            logger.error(f"Не удалось извлечь текст из файла в ZIP: {file}")
        if not extracted_text:
            raise ValueError(f"Архив {file_path} не содержит поддерживаемых файлов.")
        return "\n".join(extracted_text)
    except Exception as e:
        logger.error(f"Ошибка при обработке ZIP-архива {file_path}: {e}")
        raise Exception("Ошибка при обработке ZIP-архива.")

def extract_text_from_7z(file_path: str) -> Optional[str]:
    """Извлекает текст из файлов в 7Z-архиве с учетом MIME-типа."""
    try:
        extracted_text = []
        with py7zr.SevenZipFile(file_path, mode="r") as archive:
            archive.extractall(path="tmp")
            for root, _, files in os.walk("tmp"):
                for file in files:
                    mime_type = mimetypes.guess_type(file)[0] or "application/octet-stream"
                    file_path_extracted = os.path.join(root, file)
                    if mime_type in text_extraction_from_a_document:
                        text_extraction_function = text_extraction_from_a_document[mime_type]
                        text_content = text_extraction_function(file_path_extracted)
                        if text_content:
                            extracted_text.append(text_content)
                        else:
                            logger.error(f"Не удалось извлечь текст из файла: {file}")
        if not extracted_text:
            raise ValueError(f"7Z-архив '{file_path}' не содержит поддерживаемых файлов.")
        return "\n".join(extracted_text)
    except Exception as e:
        logger.error(f"Ошибка при обработке 7Z-архива {file_path}: {e}")
        raise Exception("Ошибка при обработке 7Z-архива.")

def extract_text_from_markdown(file_path: str) -> Optional[str]:
    """Извлекает текст из Markdown-файла."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
            if not content.strip():
                raise ValueError(f"Файл {file_path} пустой или некорректен.")
        return content
    except Exception as e:
        logger.error(f"Ошибка при обработке Markdown-файла {file_path}: {e}")
        raise Exception("Ошибка при обработке Markdown-файла.")

text_extraction_from_a_document = {
    "application/pdf": extract_text_from_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_text_from_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_text_from_excel,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_text_from_presentation,
    "application/zip": extract_text_from_zip,
    "application/x-7z-compressed": extract_text_from_7z,
    "text/markdown": extract_text_from_markdown,
}